import io
import re
import zipfile

OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def detect_format(raw: bytes) -> str:
    if raw[:4] == b"%PDF":
        return "pdf"
    if raw[:4] == b"PK\x03\x04":
        return _detect_ooxml(raw)
    if raw[:8] == OLE_MAGIC:
        return "xls"
    return "text"


def _detect_ooxml(raw: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            names = archive.namelist()
    except zipfile.BadZipFile:
        return "text"
    if any(name.startswith("ppt/") for name in names):
        return "pptx"
    if any(name.startswith("xl/") for name in names):
        return "xlsx"
    if any(name.startswith("word/") for name in names):
        return "docx"
    return "docx"


def extract_text(raw: bytes) -> tuple[str, str, int]:
    fmt = detect_format(raw)
    if fmt == "pdf":
        return _pdf_text(raw)
    if fmt == "docx":
        return _docx_text(raw), "docx", 1
    if fmt == "pptx":
        return _pptx_text(raw)
    if fmt == "xlsx":
        return _xlsx_text(raw)
    if fmt == "xls":
        return _xls_text(raw)
    return raw.decode("utf-8", errors="ignore"), "text", 1


def _pdf_text(raw: bytes) -> tuple[str, str, int]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(parts), "pdf", len(reader.pages)


def _docx_text(raw: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(raw))
    return "\n\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())


def _pptx_text(raw: bytes) -> tuple[str, str, int]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append("\t".join(cells))
    return "\n\n".join(parts), "pptx", len(presentation.slides)


def _xlsx_text(raw: bytes) -> tuple[str, str, int]:
    import openpyxl

    workbook = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    sheets = workbook.worksheets
    parts: list[str] = []
    for sheet in sheets:
        parts.append(f"# {sheet.title}")
        parts.extend(_rows_to_lines(sheet.iter_rows(values_only=True)))
    workbook.close()
    return "\n".join(parts), "xlsx", len(sheets)


def _xls_text(raw: bytes) -> tuple[str, str, int]:
    try:
        import xlrd

        book = xlrd.open_workbook(file_contents=raw)
    except Exception:
        return raw.decode("utf-8", errors="ignore"), "text", 1
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"# {sheet.name}")
        rows = (
            tuple(sheet.cell_value(row_index, col) for col in range(sheet.ncols))
            for row_index in range(sheet.nrows)
        )
        parts.extend(_rows_to_lines(rows))
    return "\n".join(parts), "xls", book.nsheets


def _rows_to_lines(rows) -> list[str]:
    lines: list[str] = []
    for row in rows:
        cells = [str(cell) for cell in row if cell is not None and str(cell).strip()]
        if cells:
            lines.append("\t".join(cells))
    return lines


def detect_lang(text: str) -> str:
    cyrillic = len(re.findall(r"[а-яё]", text, re.IGNORECASE))
    latin = len(re.findall(r"[a-z]", text, re.IGNORECASE))
    if cyrillic == 0 and latin == 0:
        return "unknown"
    if cyrillic >= latin:
        return "ru"
    return "en"


def build_docir(document_id: str, text: str, source_format: str, pages: int) -> dict:
    blocks = []
    offset = 0
    for ordinal, paragraph in enumerate(_paragraphs(text)):
        start = text.find(paragraph, offset)
        if start < 0:
            start = offset
        end = start + len(paragraph)
        offset = end
        blocks.append(
            {
                "id": f"b{ordinal}",
                "kind": "paragraph",
                "page": 1,
                "section_path": [],
                "text": paragraph,
                "char_from": start,
                "char_to": end,
            }
        )
    return {
        "schema": "docir/1",
        "document_id": document_id,
        "version": 1,
        "lang": detect_lang(text),
        "source_format": source_format,
        "pages": pages,
        "blocks": blocks,
        "full_text": text,
    }


def _paragraphs(text: str) -> list[str]:
    chunks = [chunk.strip() for chunk in re.split(r"\n\s*\n", text)]
    result = [chunk for chunk in chunks if chunk]
    return result or ([text.strip()] if text.strip() else [])
